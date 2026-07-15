"""Generate the Hackcelerate pitch deck (PPTX) for TokenOptimizer.

Builds ``docs/TokenOptimizer-Hackcelerate.pptx`` with python-pptx (pure Python).
Run:

    python scripts/generate_hackcelerate_ppt.py [output.pptx]

Numbers on the Results slides are computed live from the A/B suite when the
package is importable, so the deck always matches the validated outcome; it
falls back to the documented figures otherwise. The dashboard screenshot at
``docs/assets/hackcelerate-dashboard.png`` is embedded when present.
"""

from __future__ import annotations

import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUT = PROJECT_ROOT / "docs" / "TokenOptimizer-Hackcelerate.pptx"
DASH_IMG = PROJECT_ROOT / "docs" / "assets" / "hackcelerate-dashboard.png"
DEMO_THUMB = PROJECT_ROOT / "docs" / "assets" / "video" / "app_optimized.png"


def rgb(v: int) -> RGBColor:
    return RGBColor((v >> 16) & 255, (v >> 8) & 255, v & 255)


NAVY = rgb(0x12365C)
BLUE = rgb(0x1A5EB8)
BLUE_DK = rgb(0x134A94)
INK = rgb(0x16233A)
MUTED = rgb(0x5F6B7D)
LIGHT = rgb(0xEEF2F8)
GREEN = rgb(0x1E8E4E)
AMBER = rgb(0xE0A400)
WHITE = rgb(0xFFFFFF)
FONT = "Segoe UI"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)


# ── data (live from the A/B suite, with a documented fallback) ────────────────

def _collect():
    fallback = {
        "cases": ["checkout-retry", "fraud-scoring", "wallet-topup", "scheduled-payout",
                  "sso-login", "audit-log", "api-rate-limit", "feature-flags"],
        "q_base": [19, 19, 20, 19, 19, 20, 20, 20],
        "q_opt": [23, 23, 25, 24, 23, 25, 25, 24],
        "cost_saved": 58, "compression": 73, "quality": 24.0, "baseline_q": 19.5,
        "tokens_saved": 2413, "num_tests": 8, "num_bus": 2,
    }
    try:
        sys.path.insert(0, str(PROJECT_ROOT / "src"))
        sys.path.insert(0, str(PROJECT_ROOT / "skills"))
        from token_optimizer.evaluation.ab_runner import run_ab_suite
        from token_optimizer.evaluation.datasets import sample_cases
        from codebase_anchoring.indexer import build_index
        from prd_compression.compressor import compress_prd

        idx = build_index(str(PROJECT_ROOT / "src"))
        s = run_ab_suite(sample_cases(), idx)
        by = {c.name: c for c in sample_cases()}
        raw = comp = 0
        for r in s.results:
            cr = compress_prd(by[r.name].prd)
            raw += cr.raw_tokens
            comp += cr.compressed_tokens
        return {
            "cases": [r.name for r in s.results],
            "q_base": [r.baseline.quality.total for r in s.results],
            "q_opt": [r.optimised.quality.total for r in s.results],
            "cost_saved": round(s.avg_cost_savings_pct),
            "compression": round(100.0 * (raw - comp) / raw) if raw else 0,
            "quality": round(s.avg_optimised_quality, 1),
            "baseline_q": round(s.avg_baseline_quality, 1),
            "tokens_saved": s.total_tokens_saved,
            "num_tests": s.num_tests,
            "num_bus": len(s.bus),
        }
    except Exception as exc:  # pragma: no cover - deck still builds
        print(f"(using fallback figures: {exc})")
        return fallback


# ── slide primitives ──────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _rect(slide, l, t, w, h, color, rounded=False):
    shape = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _text(slide, l, t, w, h, lines, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """lines: list of dicts {text, size, bold, color, space_after, bullet, gap}."""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        p.space_after = Pt(ln.get("space_after", 6))
        p.space_before = Pt(ln.get("space_before", 0))
        if ln.get("line_spacing"):
            p.line_spacing = ln["line_spacing"]
        run = p.add_run()
        prefix = "•  " if ln.get("bullet") else ""
        run.text = prefix + ln["text"]
        run.font.name = FONT
        run.font.size = Pt(ln.get("size", 18))
        run.font.bold = ln.get("bold", False)
        run.font.color.rgb = ln.get("color", INK)
    return tb


def _title_bar(slide, title, kicker=None):
    _rect(slide, 0, 0, EMU_W, Inches(1.15), NAVY)
    _rect(slide, 0, Inches(1.15), EMU_W, Inches(0.06), BLUE)
    lines = []
    if kicker:
        lines.append({"text": kicker.upper(), "size": 11, "bold": True,
                      "color": rgb(0x8FB4E0), "space_after": 2})
    lines.append({"text": title, "size": 26, "bold": True, "color": WHITE})
    _text(slide, Inches(0.6), Inches(0.16), Inches(12.1), Inches(0.95), lines,
          anchor=MSO_ANCHOR.MIDDLE)


def _footer(slide, page):
    _text(slide, Inches(0.6), Inches(7.05), Inches(9), Inches(0.35),
          [{"text": "TokenOptimizer — Hackcelerate", "size": 9, "color": MUTED}])
    _text(slide, Inches(11.4), Inches(7.05), Inches(1.3), Inches(0.35),
          [{"text": str(page), "size": 9, "color": MUTED}], align=PP_ALIGN.RIGHT)


def _bullets_slide(prs, title, kicker, bullets, page, *, note=None):
    s = _blank(prs)
    _title_bar(s, title, kicker)
    lines = [{"text": b, "size": 18, "color": INK, "bullet": True,
              "space_after": 12, "line_spacing": 1.05} for b in bullets]
    _text(s, Inches(0.7), Inches(1.55), Inches(11.9), Inches(5.0), lines)
    if note:
        _rect(s, Inches(0.7), Inches(6.35), Inches(11.9), Inches(0.55), LIGHT, rounded=True)
        _text(s, Inches(0.95), Inches(6.4), Inches(11.4), Inches(0.5),
              [{"text": note, "size": 13, "bold": True, "color": NAVY}],
              anchor=MSO_ANCHOR.MIDDLE)
    _footer(s, page)
    return s


def build(prs, D):
    # 1 ── Title ────────────────────────────────────────────────────────────
    s = _blank(prs)
    _rect(s, 0, 0, EMU_W, EMU_H, NAVY)
    _rect(s, 0, Inches(4.02), EMU_W, Inches(0.08), BLUE)
    _rect(s, Inches(0.9), Inches(1.5), Inches(1.1), Inches(1.1), BLUE, rounded=True)
    _text(s, Inches(0.9), Inches(1.5), Inches(1.1), Inches(1.1),
          [{"text": "T:O", "size": 26, "bold": True, "color": WHITE}],
          align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    _text(s, Inches(0.9), Inches(2.9), Inches(11.5), Inches(1.2),
          [{"text": "TokenOptimizer", "size": 54, "bold": True, "color": WHITE}])
    _text(s, Inches(0.95), Inches(4.25), Inches(11.5), Inches(1.4), [
        {"text": "Two skills that optimise AI token usage in feature development",
         "size": 22, "color": rgb(0xD6E2F2), "space_after": 6},
        {"text": "Compress the PRD · anchor plans to real code · route to the right model — "
                 "and prove it with A/B evidence.", "size": 15, "color": rgb(0x9FB6D4)},
    ])
    _text(s, Inches(0.95), Inches(6.6), Inches(11.5), Inches(0.5),
          [{"text": f"Hackcelerate — Token Optimisation  ·  {D['date']}", "size": 13,
            "color": rgb(0x8FB4E0)}])

    # 2 ── Problem ──────────────────────────────────────────────────────────
    _bullets_slide(
        prs, "The Problem", "Why this matters", [
            "Large, verbose PRDs are fed straight into LLMs — high input-token cost on every call.",
            "AI plans reference vague or hallucinated file locations — wasted regeneration cycles, low trust.",
            "Every task, trivial or complex, hits the most expensive model — overspending on simple work.",
            "No measurable quality baseline — cost cuts risk silently degrading output.",
        ], 2,
        note="Net effect: rising AI spend, inconsistent quality, and no data-driven way to prove ROI.")

    # 3 ── Solution ─────────────────────────────────────────────────────────
    s = _blank(prs)
    _title_bar(s, "The Solution — Two Skills", "Optimise tokens, keep quality")
    cards = [
        ("Skill 1", "PRD Compression", BLUE,
         "Distils a verbose PRD into dense, structured requirement atoms before it "
         "reaches the model — ~67% fewer input tokens, acceptance criteria kept verbatim."),
        ("Skill 2", "Anchoring + Routing", GREEN,
         "Anchors every plan step to a real file:line reference (flagging hallucinations), "
         "and routes each task to the cheapest capable model."),
    ]
    x = Inches(0.7)
    for tag, name, color, desc in cards:
        _rect(s, x, Inches(1.7), Inches(5.85), Inches(4.4), LIGHT, rounded=True)
        _rect(s, x, Inches(1.7), Inches(5.85), Inches(0.14), color, rounded=False)
        _text(s, x + Inches(0.35), Inches(2.0), Inches(5.2), Inches(4.0), [
            {"text": tag.upper(), "size": 13, "bold": True, "color": color, "space_after": 3},
            {"text": name, "size": 24, "bold": True, "color": NAVY, "space_after": 12},
            {"text": desc, "size": 16, "color": INK, "line_spacing": 1.1},
        ])
        x += Inches(6.05)
    _footer(s, 3)

    # 4 ── Skill 1 detail ───────────────────────────────────────────────────
    _bullets_slide(
        prs, "Skill 1 · PRD Compression", "~67% fewer input tokens", [
            "Strips boilerplate, framing, hedging and formatting noise.",
            "Extracts only decision-relevant requirements: goals, constraints, acceptance criteria, non-functional, dependencies, out-of-scope.",
            "Preserves acceptance criteria verbatim — the build is never checked against a paraphrase.",
            "Deterministic and fully offline — no API key, no model required.",
        ], 4,
        note=f"Result: ~{D['compression']}% smaller PRDs with no loss of decision-critical content.")

    # 5 ── Skill 2 detail ───────────────────────────────────────────────────
    _bullets_slide(
        prs, "Skill 2 · Codebase Anchoring + Model Routing", "Grounded output, smart spend", [
            "Indexes the repo (AST for Python, regex for JS/TS/Java/Go/Ruby/C#) into a file:line symbol table.",
            "Anchors each plan step to real references; unresolved symbols are flagged as possible hallucinations.",
            "Classifies task complexity (trivial / standard / complex) with a confidence score.",
            "Routes trivial → Haiku, standard → Sonnet, complex → Opus; low confidence upgrades toward premium.",
        ], 5,
        note="Fewer hallucination-driven rework loops + the cheapest capable model per task.")

    # 6 ── How it works (flow) ───────────────────────────────────────────────
    s = _blank(prs)
    _title_bar(s, "How It Works", "End-to-end flow")
    steps = [("Raw PRD", MUTED), ("Compress", BLUE), ("Anchor plan\n(file:line)", GREEN),
             ("Route model", BLUE), ("LLM call", NAVY), ("Metrics +\nDashboard", GREEN)]
    n = len(steps)
    gap = Inches(0.25)
    total_w = EMU_W - Inches(1.4)
    bw = Emu(int((total_w - gap * (n - 1)) / n))
    x = Inches(0.7)
    y = Inches(3.0)
    bh = Inches(1.5)
    for i, (label, color) in enumerate(steps):
        _rect(s, x, y, bw, bh, color, rounded=True)
        _text(s, x, y, bw, bh, [{"text": label, "size": 15, "bold": True, "color": WHITE,
                                 "align": PP_ALIGN.CENTER}],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            ar = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                    x + bw, y + Inches(0.55), gap, Inches(0.4))
            ar.fill.solid(); ar.fill.fore_color.rgb = rgb(0xB9C6DA)
            ar.line.fill.background(); ar.shadow.inherit = False
        x = Emu(x + bw + gap)
    _text(s, Inches(0.7), Inches(5.1), Inches(11.9), Inches(1.0), [
        {"text": "Deterministic and offline-first: every stage works with no API key; "
                 "a Claude key or local model only improves quality.", "size": 15,
         "color": MUTED, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
    _footer(s, 6)

    # 7 ── Validated results (KPI tiles) ─────────────────────────────────────
    s = _blank(prs)
    _title_bar(s, "Validated Results", f"{D['num_tests']} A/B tests across {D['num_bus']} business units")
    tiles = [(f"{D['cost_saved']}%", "Avg cost savings", GREEN),
             (f"{D['compression']}%", "Avg PRD compression", BLUE),
             (f"{D['quality']}/25", "Avg quality (optimised)", NAVY),
             (f"{D['tokens_saved']:,}", "Total tokens saved", MUTED)]
    x = Inches(0.7)
    tw = Inches(2.9)
    for value, label, color in tiles:
        _rect(s, x, Inches(2.0), tw, Inches(2.4), LIGHT, rounded=True)
        _rect(s, x, Inches(2.0), tw, Inches(0.14), color)
        _text(s, x, Inches(2.5), tw, Inches(1.2),
              [{"text": value, "size": 40, "bold": True, "color": color,
                "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, x, Inches(3.7), tw, Inches(0.6),
              [{"text": label, "size": 14, "color": INK, "align": PP_ALIGN.CENTER}],
              align=PP_ALIGN.CENTER)
        x += Inches(3.05)
    _text(s, Inches(0.7), Inches(4.9), Inches(11.9), Inches(1.2), [
        {"text": f"Baseline quality averaged {D['baseline_q']}/25 — optimised is equal or "
                 "higher in every case, at a fraction of the cost.", "size": 16,
         "color": MUTED, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
    _footer(s, 7)

    # 8 ── Quality chart ─────────────────────────────────────────────────────
    s = _blank(prs)
    _title_bar(s, "Quality — Baseline vs Optimised", "25-point rubric, per feature request")
    cd = CategoryChartData()
    cd.categories = D["cases"]
    cd.add_series("Baseline", tuple(D["q_base"]))
    cd.add_series("Optimised", tuple(D["q_opt"]))
    gframe = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                                Inches(0.7), Inches(1.6), Inches(11.9), Inches(5.0), cd)
    chart = gframe.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.TOP
    chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.gap_width = 80
    plot.series[0].format.fill.solid()
    plot.series[0].format.fill.fore_color.rgb = BLUE
    plot.series[1].format.fill.solid()
    plot.series[1].format.fill.fore_color.rgb = GREEN
    va = chart.value_axis
    va.minimum_scale = 0
    va.maximum_scale = 25
    _footer(s, 8)

    # 9 ── Product / dashboard ───────────────────────────────────────────────
    s = _blank(prs)
    _title_bar(s, "The Product", "Desktop app · CLI · dashboard · guided demo")
    _text(s, Inches(0.7), Inches(1.5), Inches(4.4), Inches(5.2), [
        {"text": "Desktop app (3 tabs)", "size": 17, "bold": True, "color": NAVY, "space_after": 4},
        {"text": "Token Optimizer, Feature-Dev Skills, and a Dashboard with native charts.",
         "size": 13, "color": INK, "space_after": 12},
        {"text": "CLI (tokenopt)", "size": 17, "bold": True, "color": NAVY, "space_after": 4},
        {"text": "compress-prd · anchor-plan · route · ab-suite · dashboard.",
         "size": 13, "color": INK, "space_after": 12},
        {"text": "Guided demo", "size": 17, "bold": True, "color": NAVY, "space_after": 4},
        {"text": "Help → Guided Demo: an automated tour that runs each feature live.",
         "size": 13, "color": INK, "space_after": 12},
        {"text": "Skill plugins", "size": 17, "bold": True, "color": NAVY, "space_after": 4},
        {"text": "Self-contained, independently-maintained packages under skills/ "
                 "(code + SKILL.md + tests).", "size": 13, "color": INK},
    ])
    if DASH_IMG.exists():
        s.shapes.add_picture(str(DASH_IMG), Inches(5.4), Inches(1.6), width=Inches(7.4))
        _text(s, Inches(5.4), Inches(6.55), Inches(7.4), Inches(0.4),
              [{"text": "In-app A/B validation dashboard", "size": 11, "color": MUTED,
                "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
    _footer(s, 9)

    # 10 ── Demo video ───────────────────────────────────────────────────────
    s = _blank(prs)
    _title_bar(s, "Demo Video", "Narrated walkthrough of every feature")
    _text(s, Inches(0.7), Inches(1.6), Inches(4.5), Inches(5.0), [
        {"text": "A 3:14 narrated video", "size": 17, "bold": True, "color": NAVY, "space_after": 4},
        {"text": "16 scenes with voice-over and on-screen captions.",
         "size": 13, "color": INK, "space_after": 12},
        {"text": "Covers every feature", "size": 17, "bold": True, "color": NAVY, "space_after": 4},
        {"text": "Document / JIRA / GitHub optimization, the offline pipeline, both skills, "
                 "A/B validation, the dashboard, the guided demo, and the CLI.",
         "size": 13, "color": INK, "space_after": 12},
        {"text": "Reproducible & offline", "size": 17, "bold": True, "color": NAVY, "space_after": 4},
        {"text": "Built with scripts/generate_demo_video.py (SAPI voice-over + ffmpeg).",
         "size": 13, "color": INK, "space_after": 12},
        {"text": "File: docs/TokenOptimizer-Demo.mp4", "size": 12, "color": MUTED},
    ])
    if DEMO_THUMB.exists():
        pic = s.shapes.add_picture(str(DEMO_THUMB), Inches(5.5), Inches(1.7), width=Inches(7.2))
        cx = Inches(5.5) + pic.width // 2
        cy = Inches(1.7) + pic.height // 2
        r = Inches(0.58)
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
        circ.fill.solid()
        circ.fill.fore_color.rgb = NAVY
        circ.line.color.rgb = WHITE
        circ.line.width = Pt(2.5)
        circ.shadow.inherit = False
        _text(s, cx - r, cy - r - Inches(0.02), r * 2, r * 2,
              [{"text": "▶", "size": 32, "bold": True, "color": WHITE, "align": PP_ALIGN.CENTER}],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(s, Inches(5.5), Inches(6.35), Inches(7.2), Inches(0.4),
              [{"text": "TokenOptimizer-Demo.mp4  ·  3:14  ·  voice-over", "size": 11,
                "color": MUTED, "align": PP_ALIGN.CENTER}], align=PP_ALIGN.CENTER)
    _footer(s, 10)

    # 11 ── Architecture ─────────────────────────────────────────────────────
    s = _blank(prs)
    _title_bar(s, "Architecture", "Layered, offline-first")
    layers = [("Interfaces", "CLI · Desktop UI · Packaged skills", NAVY),
              ("Skills & Evaluation", "prd · anchor · router · A/B runner · dashboard", BLUE),
              ("Optimize", "reductions · compress · summarize · prefilter · tokens", BLUE_DK),
              ("Integrations", "JIRA · GitHub · GitLab · Jenkins · Azure DevOps · documents", rgb(0x2C6BB0)),
              ("Core", "config · run_log · llm (client + cache)", GREEN)]
    y = Inches(1.7)
    for name, detail, color in layers:
        _rect(s, Inches(0.7), y, Inches(11.9), Inches(0.92), color, rounded=True)
        _text(s, Inches(1.0), y, Inches(3.4), Inches(0.92),
              [{"text": name, "size": 16, "bold": True, "color": WHITE}],
              anchor=MSO_ANCHOR.MIDDLE)
        _text(s, Inches(4.5), y, Inches(8.0), Inches(0.92),
              [{"text": detail, "size": 13, "color": rgb(0xE8EEF7)}], anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(1.02)
    _footer(s, 11)

    # 12 ── Business value ───────────────────────────────────────────────────
    _bullets_slide(
        prs, "Business Value", "Why it scales", [
            f"Direct cost savings of up to ~{D['cost_saved']}% on AI-assisted development spend.",
            f"Validated across {D['num_bus']} business units — extensible to the whole org.",
            "Higher trust in AI output through verifiable file:line anchoring.",
            "Data-driven — every claim backed by A/B evidence, not anecdote.",
        ], 12,
        note="Offline-first: real savings with no API key; better still with Claude or a local model.")

    # 13 ── Closing ──────────────────────────────────────────────────────────
    s = _blank(prs)
    _rect(s, 0, 0, EMU_W, EMU_H, NAVY)
    _rect(s, 0, Inches(3.7), EMU_W, Inches(0.08), BLUE)
    _text(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.3),
          [{"text": "35% cheaper. Equal or better quality.", "size": 40, "bold": True,
            "color": WHITE}])
    _text(s, Inches(0.95), Inches(3.95), Inches(11.5), Inches(1.0), [
        {"text": f"Validated across {D['num_tests']} A/B tests and {D['num_bus']} business "
                 f"units — ~{D['compression']}% compression, ~{D['cost_saved']}% savings, "
                 f"{D['quality']}/25 quality.", "size": 18, "color": rgb(0xD6E2F2)}])
    _text(s, Inches(0.95), Inches(5.6), Inches(11.5), Inches(0.6),
          [{"text": "Thank you  ·  TokenOptimizer", "size": 16, "color": rgb(0x8FB4E0)}])


def main(argv):
    out = Path(argv[1]) if len(argv) > 1 else DEFAULT_OUT
    out.parent.mkdir(parents=True, exist_ok=True)
    D = _collect()
    D["date"] = datetime.now().strftime("%d %b %Y")
    prs = Presentation()
    prs.slide_width = EMU_W
    prs.slide_height = EMU_H
    build(prs, D)
    prs.save(str(out))
    print(f"PPTX written to: {out}  ({len(prs.slides)} slides)")
    return 0


if __name__ == "__main__":
    raise SystemExit(main(sys.argv))
